"""
PDF text extraction and section detection for research papers.
Uses PyMuPDF (fitz) with pdfminer as fallback.
"""

import re
import io
from typing import Optional


def extract_text(file_bytes: bytes) -> tuple:
    """
    Extract all text from a PDF file.
    Returns (text, error). Error is empty string on success.
    """
    # Try PyMuPDF first (fastest, best quality)
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        text = "\n".join(pages)
        if text.strip():
            return text, ""
    except ImportError:
        pass
    except Exception as e:
        return "", f"PyMuPDF error: {e}"

    # Fallback: pdfminer
    try:
        from pdfminer.high_level import extract_text_to_fp
        from pdfminer.layout import LAParams
        out = io.StringIO()
        extract_text_to_fp(io.BytesIO(file_bytes), out, laparams=LAParams())
        text = out.getvalue()
        if text.strip():
            return text, ""
    except ImportError:
        pass
    except Exception as e:
        return "", f"pdfminer error: {e}"

    return "", (
        "No PDF extraction library available. "
        "Install PyMuPDF: pip install PyMuPDF"
    )


def _find_section(text_lower: str, text: str, keywords: list, max_len: int = 1500) -> str:
    """Find the first matching section keyword and return its content."""
    for kw in keywords:
        idx = text_lower.find(kw)
        if idx >= 0:
            # Find a reasonable end (next major section or max_len chars)
            snippet = text[idx: idx + max_len]
            return snippet.strip()
    return ""


def extract_sections(text: str) -> dict:
    """
    Heuristically split a research paper into its canonical sections.
    Returns a dict with keys: abstract, introduction, methods, results, conclusions.
    """
    tl = text.lower()
    sections = {}

    abstract = _find_section(tl, text,
        ["abstract\n", "abstract—", "abstract:", "abstract "],
        max_len=1200)
    if abstract:
        sections["abstract"] = abstract

    intro = _find_section(tl, text,
        ["1. introduction", "1introduction", "introduction\n", "introduction\r"],
        max_len=1500)
    if intro:
        sections["introduction"] = intro

    methods = _find_section(tl, text,
        ["materials and methods", "methodology\n", "methods\n", "2. method", "3. method",
         "experimental methods", "study design"],
        max_len=1500)
    if methods:
        sections["methods"] = methods

    results = _find_section(tl, text,
        ["results\n", "results and discussion", "findings\n", "3. result", "4. result"],
        max_len=1500)
    if results:
        sections["results"] = results

    conclusions = _find_section(tl, text,
        ["conclusion\n", "conclusions\n", "discussion and conclusion",
         "summary and conclusion", "closing remarks"],
        max_len=1000)
    if conclusions:
        sections["conclusions"] = conclusions

    return sections


def extract_metadata(text: str) -> dict:
    """Try to pull title, authors and DOI from the first page."""
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    title = lines[0][:200] if lines else ""
    doi_match = re.search(r"10\.\d{4,9}/[^\s]+", text)
    doi = doi_match.group(0) if doi_match else ""
    return {"title": title, "doi": doi}


def build_ai_prompt(sections: dict, max_chars: int = 3000) -> str:
    """Build a condensed prompt string from extracted sections."""
    parts = []
    for name, content in sections.items():
        parts.append(f"## {name.upper()}\n{content[:600]}")
    return "\n\n".join(parts)[:max_chars]


# ── Slide generation ──────────────────────────────────────────────────────────

def generate_slides(title: str, sections: dict, ai_summary: Optional[dict] = None) -> bytes:
    """
    Generate a PowerPoint (.pptx) from paper sections and optional AI summary.
    Returns the file bytes, or empty bytes if python-pptx is not installed.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        def _add_title_slide(prs, title):
            layout = prs.slide_layouts[0]
            slide  = prs.slides.add_slide(layout)
            slide.shapes.title.text = title[:200]
            try:
                slide.placeholders[1].text = "Auto-generated Research Summary · ORC Dashboard"
            except Exception:
                pass
            return slide

        def _add_content_slide(prs, heading, body_text):
            layout = prs.slide_layouts[1]
            slide  = prs.slides.add_slide(layout)
            slide.shapes.title.text = heading
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            tf.text = body_text[:600]
            return slide

        _add_title_slide(prs, title)

        section_headings = {
            "abstract":     "Abstract",
            "introduction": "Introduction",
            "methods":      "Methodology",
            "results":      "Results",
            "conclusions":  "Conclusions",
        }

        if ai_summary:
            for key, heading in [
                ("overview",    "Overview"),
                ("objectives",  "Objectives"),
                ("methods",     "Methodology"),
                ("results",     "Key Results"),
                ("conclusion",  "Conclusion"),
            ]:
                val = ai_summary.get(key)
                if not val:
                    continue
                body = "\n• ".join(val) if isinstance(val, list) else str(val)
                _add_content_slide(prs, heading, body)
        else:
            for key, heading in section_headings.items():
                if key in sections:
                    _add_content_slide(prs, heading, sections[key][:600])

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    except ImportError:
        return b""
    except Exception:
        return b""
